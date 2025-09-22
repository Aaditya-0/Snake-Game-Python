from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

# Slide 1: Problem Statement
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Snake Game - Problem Statement"
slide.placeholders[1].text = (
    "Design and implement a classic Snake Game using Python and Pygame.\n"
    "The game should allow the player to control a snake, eat food, grow in length, and avoid collisions with walls or itself."
)

# Slide 2: Requirements & Tech Stack
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Requirements & Tech Stack"
slide.placeholders[1].text = (
    "Requirements:\n"
    "- Player controls snake with arrow keys\n"
    "- Snake grows when eating food\n"
    "- Game ends on collision\n"
    "- Score tracking\n\n"
    "Tech Stack:\n"
    "- Python 3.x\n"
    "- Pygame library\n"
    "- python-pptx (for this presentation)"
)

# Slide 3: Module & File Structure
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Module & File Structure"
slide.placeholders[1].text = (
    "Project Structure:\n"
    "src/\n"
    "  └── main.py (game logic)\n"
    "  └── resources/\n"
    "      ├── snake.png\n"
    "      ├── block.png\n"
    "      └── food.jpeg\n"
    "\nmain.py Modules:\n"
    "- SnakeGame class\n"
    "- Game loop (run)\n"
    "- Drawing & event handling"
)

# Slide 4: Backend Logic - Game Loop
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Backend Logic: Game Loop"
slide.placeholders[1].text = (
    "The main game loop handles:\n"
    "- Processing user input (arrow keys)\n"
    "- Moving the snake\n"
    "- Checking for collisions (walls, self)\n"
    "- Eating food and growing\n"
    "- Drawing background, snake, and food\n"
    "- Updating the display"
)

# Slide 5: Backend Logic - Collision & Food
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Backend Logic: Collision & Food"
slide.placeholders[1].text = (
    "Collision Detection:\n"
    "- Wall collision: snake head out of bounds\n"
    "- Self collision: snake head in body\n"
    "- On collision: show game over screen\n\n"
    "Food Handling:\n"
    "- Randomly place food not on snake\n"
    "- On eating: increase score, grow snake"
)

# Slide 6: Challenges Faced
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Challenges Faced"
slide.placeholders[1].text = (
    "- Handling fast user input and preventing reverse direction\n"
    "- Ensuring food does not spawn on the snake\n"
    "- Smooth rendering and frame rate control\n"
    "- Managing game restart and state transitions"
)

# Slide 7: Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Conclusion"
slide.placeholders[1].text = (
    "The Snake Game project demonstrates:\n"
    "- Use of Pygame for 2D game development\n"
    "- Event-driven programming in Python\n"
    "- Basic game logic and collision handling\n"
    "- Modular code structure for maintainability"
)

prs.save("Snake_Game_Presentation.pptx")
print("Presentation created: Snake_Game_Presentation.pptx")